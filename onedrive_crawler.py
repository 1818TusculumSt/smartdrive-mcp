import os
import msal
import requests
import io
import json
from pathlib import Path
from dotenv import load_dotenv
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import olefile
from PIL import Image
import warnings
import csv
import zipfile

# Suppress PyTorch warnings about GPU/accelerator
warnings.filterwarnings('ignore', category=UserWarning, module='torch')

import easyocr
from pinecone import Pinecone
from embeddings import EmbeddingProvider
from config import settings

load_dotenv()

CLIENT_ID = os.getenv("MICROSOFT_CLIENT_ID")
TENANT_ID = os.getenv("MICROSOFT_TENANT_ID")
AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPES = ["Files.Read.All", "User.Read"]

# Token cache file
TOKEN_CACHE_FILE = Path.home() / ".smartdrive_token_cache.json"

# Folder skip cache file (remembers your choices)
FOLDER_SKIP_CACHE_FILE = Path.home() / ".smartdrive_folder_skip_cache.json"

# Global settings for this run
EXTRACT_ZIP_CONTENTS = False  # Will be set based on user choice

# Initialize Pinecone
pc = Pinecone(api_key=settings.PINECONE_API_KEY)
index = pc.Index(
    name=settings.PINECONE_INDEX_NAME,
    host=settings.PINECONE_HOST
)

# Initialize embedding provider
# Show correct model name based on provider
if settings.EMBEDDING_PROVIDER == "voyage":
    model_display = settings.VOYAGE_MODEL
else:
    model_display = settings.EMBEDDING_MODEL

print(f"🧠 Loading {settings.EMBEDDING_PROVIDER} embedding provider ({model_display})...")
embedding_provider = EmbeddingProvider()
print("✅ Embedding provider loaded\n")

# Initialize EasyOCR reader (lazy-loaded on first use)
ocr_reader = None

# Azure Computer Vision configuration
AZURE_VISION_KEY = os.getenv("AZURE_VISION_KEY")
AZURE_VISION_ENDPOINT = os.getenv("AZURE_VISION_ENDPOINT")
USE_AZURE_OCR = bool(AZURE_VISION_KEY and AZURE_VISION_ENDPOINT)

if USE_AZURE_OCR:
    print("☁️  Azure Computer Vision OCR enabled (10-20x faster than local!)")
    from azure.ai.vision.imageanalysis import ImageAnalysisClient
    from azure.ai.vision.imageanalysis.models import VisualFeatures
    from azure.core.credentials import AzureKeyCredential

    azure_client = ImageAnalysisClient(
        endpoint=AZURE_VISION_ENDPOINT,
        credential=AzureKeyCredential(AZURE_VISION_KEY)
    )
    print("✅ Azure OCR ready\n")
else:
    print("💻 Using local EasyOCR (Azure OCR not configured)")
    print("   💡 Tip: Add AZURE_VISION_KEY to .env for 10-20x faster OCR!\n")

def get_ocr_reader():
    """Lazy-load EasyOCR reader (downloads models on first use)"""
    global ocr_reader
    if ocr_reader is None:
        print("🔍 Loading local OCR model (first time only, may take a moment)...")
        ocr_reader = easyocr.Reader(['en'], gpu=False, verbose=False)
        print("✅ Local OCR model loaded\n")
    return ocr_reader

def ocr_image_with_azure(image_data):
    """Use Azure Computer Vision to extract text from image

    Args:
        image_data: Image bytes

    Returns:
        Extracted text or None if failed
    """
    try:
        result = azure_client.analyze(
            image_data=image_data,
            visual_features=[VisualFeatures.READ]
        )

        if result.read is None or result.read.blocks is None:
            return None

        # Extract text from all blocks
        text_parts = []
        for block in result.read.blocks:
            for line in block.lines:
                text_parts.append(line.text)

        return "\n".join(text_parts)
    except Exception as e:
        print(f"      ⚠️ Azure OCR failed: {e}")
        return None

def load_token_cache():
    """Load token cache from file"""
    cache = msal.SerializableTokenCache()
    if TOKEN_CACHE_FILE.exists():
        with open(TOKEN_CACHE_FILE, 'r') as f:
            cache.deserialize(f.read())
    return cache

def save_token_cache(cache):
    """Save token cache to file"""
    if cache.has_state_changed:
        with open(TOKEN_CACHE_FILE, 'w') as f:
            f.write(cache.serialize())

def get_access_token(silent_only=False):
    """Authenticate using device code flow with token caching

    Args:
        silent_only: If True, only attempt silent token refresh (no interactive prompts)
    """
    cache = load_token_cache()
    app = msal.PublicClientApplication(CLIENT_ID, authority=AUTHORITY, token_cache=cache)

    # Try to get token silently from cache first
    accounts = app.get_accounts()
    if accounts:
        if not silent_only:
            print("🔄 Using cached credentials...")
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            save_token_cache(cache)
            if not silent_only:
                print("✅ Authentication successful (cached)\n")
            return result["access_token"]

    # If silent_only mode, return None (caller will handle)
    if silent_only:
        return None

    # If silent auth fails, do interactive device flow
    print("🔐 No cached credentials found, starting authentication...")
    flow = app.initiate_device_flow(scopes=SCOPES)

    if "user_code" not in flow:
        raise Exception(f"Failed to create device flow: {flow.get('error_description', 'Unknown error')}")

    print(f"\n🔐 Go to: {flow['verification_uri']}")
    print(f"📱 Enter code: {flow['user_code']}\n")

    result = app.acquire_token_by_device_flow(flow)

    if "access_token" in result:
        save_token_cache(cache)
        return result["access_token"]
    else:
        raise Exception(f"Authentication failed: {result.get('error_description')}")

def extract_text_from_zip_item(file_name, content):
    """Extract text from files inside zip archives"""
    file_name_lower = file_name.lower()

    try:
        # PDF
        if file_name_lower.endswith('.pdf'):
            pdf = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page_num in range(min(pdf.page_count, 10)):  # Limit to 10 pages per PDF in zip
                text += pdf[page_num].get_text() + "\n"
            pdf.close()
            return text.strip()

        # Word (.docx)
        elif file_name_lower.endswith('.docx'):
            doc = Document(io.BytesIO(content))
            text = "\n".join([para.text for para in doc.paragraphs])
            return text.strip()

        # PowerPoint (.pptx)
        elif file_name_lower.endswith('.pptx'):
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n".join(text_parts).strip()

        # Markdown
        elif file_name_lower.endswith(('.md', '.markdown')):
            return content.decode('utf-8', errors='ignore').strip()

        # Plain text
        elif file_name_lower.endswith('.txt'):
            return content.decode('utf-8', errors='ignore').strip()

        # CSV
        elif file_name_lower.endswith('.csv'):
            text_content = content.decode('utf-8', errors='ignore')
            csv_reader = csv.reader(io.StringIO(text_content))
            rows = list(csv_reader)[:50]  # Limit rows
            text_parts = []
            for row in rows:
                row_text = " | ".join([str(cell) for cell in row if cell])
                if row_text.strip():
                    text_parts.append(row_text)
            return "\n".join(text_parts).strip()

        # JSON
        elif file_name_lower.endswith('.json'):
            try:
                text_content = content.decode('utf-8', errors='ignore')
                json_data = json.loads(text_content)
                formatted_json = json.dumps(json_data, indent=2, ensure_ascii=False)
                return formatted_json[:5000]  # Limit JSON size
            except:
                return content.decode('utf-8', errors='ignore').strip()[:5000]

        # Excel
        elif file_name_lower.endswith(('.xlsx', '.xlsm')):
            workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
            text_parts = []
            for sheet_name in workbook.sheetnames[:3]:  # Limit to 3 sheets
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                for idx, row in enumerate(sheet.iter_rows(values_only=True)):
                    if idx > 50:  # Limit rows per sheet
                        break
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
            workbook.close()
            return "\n".join(text_parts).strip()

        else:
            return None

    except Exception as e:
        return None

def extract_text_from_file(token, file_item):
    """Download and extract text from supported file types"""
    headers = {"Authorization": f"Bearer {token}"}
    download_url = file_item.get("@microsoft.graph.downloadUrl")
    
    if not download_url:
        return None

    # Download file content with retry logic
    max_retries = 3
    for attempt in range(max_retries):
        try:
            response = requests.get(download_url, timeout=60)  # 60 second timeout
            if response.status_code == 200:
                break
            elif attempt < max_retries - 1:
                print(f"   ⚠️ Download failed (HTTP {response.status_code}), retrying... ({attempt + 1}/{max_retries})")
                continue
            else:
                return None
        except (requests.exceptions.ConnectionError, requests.exceptions.ChunkedEncodingError) as e:
            if attempt < max_retries - 1:
                wait_time = (attempt + 1) * 2  # Exponential backoff: 2s, 4s, 6s
                print(f"   ⚠️ Connection error, retrying in {wait_time}s... ({attempt + 1}/{max_retries})")
                import time
                time.sleep(wait_time)
                continue
            else:
                print(f"   ❌ Download failed after {max_retries} attempts: {type(e).__name__}")
                return None
        except requests.exceptions.Timeout:
            if attempt < max_retries - 1:
                print(f"   ⏱️ Download timeout, retrying... ({attempt + 1}/{max_retries})")
                continue
            else:
                print(f"   ❌ Download timeout after {max_retries} attempts")
                return None
        except Exception as e:
            print(f"   ❌ Unexpected download error: {e}")
            return None
    
    file_name = file_item["name"].lower()
    content = response.content
    
    try:
        # PDF extraction (with OCR fallback for scanned docs)
        if file_name.endswith('.pdf'):
            pdf = fitz.open(stream=content, filetype="pdf")
            text = ""

            # Try normal text extraction first
            for page_num in range(pdf.page_count):
                page = pdf[page_num]
                page_text = page.get_text()
                if page_text:
                    text += page_text + "\n"

            # If no text found or very little text (likely scanned), use OCR
            if len(text.strip()) < 50:
                print(f"   🔍 Scanned PDF detected ({pdf.page_count} pages), using OCR...")

                azure_ocr_succeeded = False
                if USE_AZURE_OCR:
                    # Use Azure Computer Vision OCR (10-20x faster!)
                    print(f"      ☁️  Using Azure OCR (1-3 seconds per page)...")
                    try:
                        ocr_text = ""
                        for page_num in range(pdf.page_count):
                            print(f"      📄 Page {page_num+1}/{pdf.page_count}...", end=" ", flush=True)
                            page = pdf[page_num]
                            # Render page to image
                            pix = page.get_pixmap(dpi=300)
                            # Convert to bytes (PNG format)
                            img_bytes = pix.tobytes("png")
                            # Run Azure OCR
                            page_text = ocr_image_with_azure(img_bytes)
                            if page_text:
                                ocr_text += f"=== Page {page_num+1} ===\n{page_text}\n"
                                print("✓")
                            else:
                                print("⚠️")
                        print(f"      ✅ Azure OCR complete!")
                        text = ocr_text
                        azure_ocr_succeeded = True
                    except Exception as ocr_error:
                        print(f"      ⚠️ Azure OCR failed: {ocr_error}")
                        print(f"      🔄 Falling back to local OCR...")

                if not azure_ocr_succeeded and len(text.strip()) < 50:
                    # Use local EasyOCR (slower but works offline)
                    print(f"      💻 Using local OCR (10-30 seconds per page)...")
                    try:
                        reader = get_ocr_reader()
                        ocr_text = ""
                        for page_num in range(pdf.page_count):
                            print(f"      📄 Page {page_num+1}/{pdf.page_count}...", end=" ", flush=True)
                            page = pdf[page_num]
                            # Render page to image (pixmap) at 300 DPI
                            pix = page.get_pixmap(dpi=300)
                            # Convert pixmap to numpy array for EasyOCR
                            import numpy as np
                            img_array = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, 3)
                            # Run OCR on the image
                            result = reader.readtext(img_array, detail=0, paragraph=True)
                            page_text = "\n".join(result)
                            ocr_text += f"=== Page {page_num+1} ===\n{page_text}\n"
                            print("✓")
                        print(f"      ✅ Local OCR complete!")
                        text = ocr_text
                    except Exception as ocr_error:
                        print(f"   ⚠️ Local OCR failed: {ocr_error}")
                        # Return whatever text we got from normal extraction
                        pass

            pdf.close()
            return text.strip()

        # Word doc extraction (.docx)
        elif file_name.endswith('.docx'):
            try:
                doc = Document(io.BytesIO(content))
                text = "\n".join([para.text for para in doc.paragraphs])
                return text.strip()
            except Exception as docx_error:
                print(f"   ⚠️ Word document parsing failed: {docx_error}")
                return None

        # Legacy Word doc extraction (.doc)
        elif file_name.endswith('.doc'):
            try:
                ole = olefile.OleFileIO(content)
                if ole.exists('WordDocument'):
                    # Try to extract text from Word 97-2003 format
                    # This is a simplified extraction - won't get all formatting
                    word_stream = ole.openstream('WordDocument')
                    data = word_stream.read()
                    # Simple heuristic: extract printable ASCII/Unicode text
                    text = data.decode('latin-1', errors='ignore')
                    # Clean up - remove non-printable chars but keep newlines/spaces
                    text = ''.join(c for c in text if c.isprintable() or c in '\n\r\t ')
                    # Remove excessive whitespace
                    text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
                    ole.close()
                    return text.strip() if len(text.strip()) > 20 else None
                else:
                    print(f"   ⚠️ Not a valid .doc file")
                    ole.close()
                    return None
            except Exception as e:
                print(f"   ⚠️ Failed to parse .doc file: {e}")
                return None

        # PowerPoint extraction (.pptx)
        elif file_name.endswith('.pptx'):
            try:
                prs = Presentation(io.BytesIO(content))
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"=== Slide {slide_num} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text_parts.append(shape.text)
                return "\n".join(text_parts).strip()
            except Exception as pptx_error:
                print(f"   ⚠️ PowerPoint parsing failed: {pptx_error}")
                return None

        # Legacy PowerPoint (.ppt) - use Apache Tika
        elif file_name.endswith('.ppt'):
            print(f"   📋 Legacy .ppt format - using Apache Tika for extraction...")
            try:
                from tika import parser
                # Parse .ppt file with Tika
                parsed = parser.from_buffer(content)
                text = parsed.get("content", "")
                if text and len(text.strip()) > 20:
                    print(f"   ✅ Extracted {len(text)} characters via Tika")
                    return text.strip()
                else:
                    print(f"   ⚠️ Tika extraction returned no text")
                    return f"File: {file_name}\nType: Legacy PowerPoint (.ppt)\nNote: No text extracted"
            except ImportError:
                print(f"   ⚠️ Apache Tika not installed - install 'tika' package and Java 11+")
                print(f"   📋 Indexing filename/path only")
                return f"File: {file_name}\nType: Legacy PowerPoint (.ppt)\nNote: Tika not available"
            except Exception as tika_error:
                print(f"   ⚠️ Tika extraction failed: {tika_error}")
                print(f"   📋 Indexing filename/path only")
                return f"File: {file_name}\nType: Legacy PowerPoint (.ppt)\nNote: Extraction failed"

        # Publisher files (.pub) - index metadata only
        elif file_name.endswith('.pub'):
            print(f"   📋 Publisher file - indexing filename/path only")
            return f"File: {file_name}\nType: Microsoft Publisher (.pub)\nNote: Full-text extraction not supported"

        # Markdown extraction (.md)
        elif file_name.endswith(('.md', '.markdown')):
            try:
                # Markdown is already plain text, just decode it
                text = content.decode('utf-8', errors='ignore')
                return text.strip()
            except Exception as md_error:
                print(f"   ⚠️ Markdown parsing failed: {md_error}")
                return None

        # Excel extraction
        elif file_name.endswith(('.xlsx', '.xlsm', '.xltx', '.xltm')):
            try:
                workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
                text_parts = []
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_parts.append(f"=== Sheet: {sheet_name} ===")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            text_parts.append(row_text)
                workbook.close()
                return "\n".join(text_parts).strip()
            except Exception as excel_error:
                print(f"   ⚠️ Excel parsing failed: {excel_error}")
                return None

        # Image extraction with OCR
        elif file_name.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif')):
            azure_ocr_succeeded = False

            if USE_AZURE_OCR:
                # Try Azure OCR first (10-20x faster!)
                print(f"   ☁️  Using Azure OCR...")
                try:
                    text = ocr_image_with_azure(content)
                    if text:
                        azure_ocr_succeeded = True
                        return text.strip()
                except Exception as azure_error:
                    print(f"   ⚠️ Azure OCR failed: {azure_error}")
                    print(f"   🔄 Falling back to local OCR...")

            if not azure_ocr_succeeded:
                # Use local EasyOCR as fallback
                print(f"   💻 Using local OCR...")
                import numpy as np
                reader = get_ocr_reader()
                image = Image.open(io.BytesIO(content))
                # Convert PIL image to numpy array
                img_array = np.array(image)
                # Run OCR
                result = reader.readtext(img_array, detail=0, paragraph=True)
                text = "\n".join(result)
                return text.strip() if text.strip() else None

        # Plain text
        elif file_name.endswith('.txt'):
            return content.decode('utf-8', errors='ignore').strip()

        # CSV extraction
        elif file_name.endswith('.csv'):
            try:
                # Try UTF-8 first, fallback to latin-1
                try:
                    text_content = content.decode('utf-8')
                except UnicodeDecodeError:
                    text_content = content.decode('latin-1', errors='ignore')

                # Parse CSV
                csv_reader = csv.reader(io.StringIO(text_content))
                rows = list(csv_reader)

                # Format as readable text
                text_parts = []
                for row in rows:
                    row_text = " | ".join([str(cell) for cell in row if cell])
                    if row_text.strip():
                        text_parts.append(row_text)

                return "\n".join(text_parts).strip()
            except Exception as csv_error:
                print(f"   ⚠️ CSV parsing failed: {csv_error}")
                return None

        # JSON extraction
        elif file_name.endswith('.json'):
            try:
                text_content = content.decode('utf-8', errors='ignore')
                json_data = json.loads(text_content)

                # Pretty print JSON for readability
                formatted_json = json.dumps(json_data, indent=2, ensure_ascii=False)
                return formatted_json
            except Exception as json_error:
                print(f"   ⚠️ JSON parsing failed: {json_error}")
                # Fall back to raw text if JSON parsing fails
                try:
                    return content.decode('utf-8', errors='ignore').strip()
                except:
                    return None

        # ZIP file handling
        elif file_name.endswith('.zip'):
            try:
                zip_file = zipfile.ZipFile(io.BytesIO(content))
                file_list = zip_file.namelist()

                if EXTRACT_ZIP_CONTENTS:
                    # Option 2: Extract and index contents
                    print(f"   📦 Extracting {len(file_list)} files from archive...")
                    extracted_texts = []
                    extracted_texts.append(f"=== Archive: {file_name} ===")
                    extracted_texts.append(f"Contains {len(file_list)} files\n")

                    for zip_item_name in file_list[:50]:  # Limit to 50 files per archive
                        try:
                            # Skip directories
                            if zip_item_name.endswith('/'):
                                continue

                            zip_item_ext = zip_item_name.lower().split('.')[-1] if '.' in zip_item_name else ''
                            supported_in_zip = ['pdf', 'docx', 'doc', 'pptx', 'txt', 'csv', 'json', 'md', 'markdown', 'xlsx']

                            if zip_item_ext in supported_in_zip:
                                print(f"      📄 Extracting: {zip_item_name}")
                                zip_item_content = zip_file.read(zip_item_name)

                                # Create a fake file_item for extract_text_from_file
                                fake_item = {
                                    'name': zip_item_name,
                                    '@microsoft.graph.downloadUrl': None  # Signal we have content already
                                }

                                # Use a custom extraction for zip contents
                                zip_text = extract_text_from_zip_item(zip_item_name, zip_item_content)
                                if zip_text:
                                    extracted_texts.append(f"\n=== File: {zip_item_name} ===")
                                    extracted_texts.append(zip_text[:2000])  # Limit each file to 2000 chars
                        except Exception as zip_item_error:
                            print(f"      ⚠️ Failed to extract {zip_item_name}: {zip_item_error}")
                            continue

                    zip_file.close()
                    return "\n".join(extracted_texts).strip()
                else:
                    # Option 3: Just list contents (default)
                    text_parts = []
                    text_parts.append(f"=== Archive: {file_name} ===")
                    text_parts.append(f"Contains {len(file_list)} files:")
                    for item in file_list[:100]:  # List up to 100 files
                        text_parts.append(f"  • {item}")
                    if len(file_list) > 100:
                        text_parts.append(f"  ... and {len(file_list) - 100} more files")

                    zip_file.close()
                    return "\n".join(text_parts).strip()

            except Exception as zip_error:
                print(f"   ⚠️ ZIP processing failed: {zip_error}")
                return None

        else:
            return None

    except Exception as e:
        print(f"❌ Failed to extract {file_name}: {e}")
        return None

def generate_vector_id(file_path):
    """Generate deterministic vector ID from file path"""
    import hashlib
    # Use MD5 hash of file path for consistent, URL-safe IDs
    return hashlib.md5(file_path.encode('utf-8')).hexdigest()

def upload_to_pinecone(files_data, check_existing=True):
    """Upload extracted files to Pinecone with embeddings (incremental sync enabled)

    Args:
        files_data: List of file data dictionaries
        check_existing: If True, check Pinecone for existing files and skip unchanged ones
    """
    print(f"\n📤 Processing {len(files_data)} files for Pinecone upload...")

    # If incremental sync enabled, check which files already exist
    existing_files = {}
    if check_existing:
        print(f"   🔍 Checking Pinecone for existing files...")
        try:
            # Query Pinecone for all existing file paths in batches
            vector_ids = [generate_vector_id(file_data["path"]) for file_data in files_data]

            # Fetch in batches of 100 (Pinecone limit)
            for i in range(0, len(vector_ids), 100):
                batch_ids = vector_ids[i:i+100]
                try:
                    result = index.fetch(ids=batch_ids, namespace="smartdrive")
                    for vid, vector in result.get('vectors', {}).items():
                        # Map back to file path
                        file_path = vector.get('metadata', {}).get('file_path', '')
                        if file_path:
                            existing_files[file_path] = {
                                "modified": vector.get('metadata', {}).get('modified', ''),
                                "size": vector.get('metadata', {}).get('size', 0)
                            }
                except Exception as e:
                    # Batch fetch failed, continue
                    pass

            if existing_files:
                print(f"   ✅ Found {len(existing_files)} existing files in index")
        except Exception as e:
            print(f"   ⚠️ Could not check existing files: {e}")
            print(f"   ℹ️  Proceeding without incremental sync")

    vectors = []
    skipped_count = 0
    updated_count = 0
    new_count = 0

    for file_data in files_data:
        file_path = file_data["path"]
        file_modified = file_data["modified"]
        file_size = file_data["size"]

        # Check if file exists and hasn't changed
        if check_existing and file_path in existing_files:
            existing = existing_files[file_path]
            if existing["modified"] == file_modified and existing["size"] == file_size:
                skipped_count += 1
                print(f"   ⏭️  Skipped (unchanged): {file_data['name']}")
                continue
            else:
                updated_count += 1
                print(f"   🔄 Updating (modified): {file_data['name']}")
        else:
            new_count += 1
            print(f"   ➕ Adding (new): {file_data['name']}")

        # Skip files with too little content (won't generate useful embeddings)
        full_text = file_data["text"]

        # Filter out files with < 50 characters of actual content
        # (headers, empty CSVs, metadata-only entries)
        if len(full_text.strip()) < 50:
            print(f"      ⏭️  Skipped (insufficient content: {len(full_text.strip())} chars)")
            skipped_count += 1
            continue

        # Truncate for embedding based on provider limits
        # Voyage AI: 32K tokens = ~128K chars (entire 50+ page PDFs!)
        # Pinecone llama: 2048 tokens = ~8K chars
        # Default to 128K for maximum document coverage with Voyage
        text_for_embedding = full_text[:128000]

        # Generate embedding
        embedding = embedding_provider.get_embedding_sync(text_for_embedding)

        if embedding is None:
            print(f"      ⚠️ Failed to generate embedding")
            continue

        embedding = embedding.tolist()

        # Generate deterministic vector ID from file path
        vector_id = generate_vector_id(file_path)

        vectors.append({
            "id": vector_id,
            "values": embedding,
            "metadata": {
                "file_name": file_data["name"],
                "file_path": file_path,
                "size": file_size,
                "modified": file_modified,
                "text_preview": full_text[:2000]  # Store larger preview (first 2000 chars)
            }
        })

    # Upsert to Pinecone (upsert will update existing or insert new)
    if vectors:
        index.upsert(vectors=vectors, namespace="smartdrive")
        print(f"\n✅ Pinecone upload complete:")
        print(f"   ➕ New files: {new_count}")
        print(f"   🔄 Updated files: {updated_count}")
        print(f"   ⏭️  Skipped (unchanged): {skipped_count}")
        print(f"   📊 Total processed: {len(files_data)} files")
    else:
        print(f"\n✅ No files needed uploading (all {skipped_count} unchanged)")

def delete_folder_from_index(folder_path):
    """Delete all files from a folder (and subfolders) from Pinecone index

    Args:
        folder_path: Path to folder to delete (e.g., "/Documents/MyFolder")

    Returns:
        Number of files deleted
    """
    print(f"\n🗑️  Searching for files in: {folder_path}")

    try:
        # Query Pinecone for all vectors with metadata.file_path starting with folder_path
        # We'll need to use filter query - but Pinecone free tier doesn't support metadata filtering
        # So we'll use a different approach: list all IDs and fetch metadata

        # Get stats to see total vectors
        stats = index.describe_index_stats()
        namespace_stats = stats.get('namespaces', {}).get('smartdrive', {})
        total_vectors = namespace_stats.get('vector_count', 0)

        if total_vectors == 0:
            print(f"ℹ️  Index is empty, nothing to delete")
            return 0

        print(f"   📊 Scanning {total_vectors} vectors in index...")

        # Since we can't filter by metadata in free tier, we need to:
        # 1. Generate all possible vector IDs for files in that folder
        # 2. Or scan through all vectors (expensive but works)

        # Better approach: Keep track during discovery
        print(f"\n⚠️  Note: To delete specific folders, we need to scan your OneDrive")
        print(f"   This will discover which files are in that folder")

        # Get token
        token = get_access_token(silent_only=True)
        if not token:
            print(f"❌ Need to refresh authentication")
            token = get_access_token()

        # Find the folder in OneDrive
        headers = {"Authorization": f"Bearer {token}"}
        base_url = "https://graph.microsoft.com/v1.0/me/drive"

        # Parse folder path to find the folder
        # folder_path format: "/Documents/MyFolder/SubFolder"
        folder_url = f"{base_url}/root:{folder_path}"
        response = requests.get(folder_url, headers=headers)

        if response.status_code != 200:
            print(f"❌ Folder not found in OneDrive: {folder_path}")
            return 0

        folder_id = response.json()["id"]
        print(f"   ✅ Found folder in OneDrive")

        # Discover all files in this folder and subfolders
        print(f"   🔍 Discovering files in folder...")
        file_paths = []
        discover_files_in_folder(token, folder_id, folder_path, file_paths)

        print(f"   ✅ Found {len(file_paths)} files to delete")

        if len(file_paths) == 0:
            print(f"ℹ️  No files found in folder")
            return 0

        # Generate vector IDs for these files
        vector_ids = [generate_vector_id(fp) for fp in file_paths]

        # Delete in batches of 100
        deleted_count = 0
        for i in range(0, len(vector_ids), 100):
            batch_ids = vector_ids[i:i+100]
            try:
                index.delete(ids=batch_ids, namespace="smartdrive")
                deleted_count += len(batch_ids)
                print(f"   🗑️  Deleted batch: {deleted_count}/{len(vector_ids)} vectors")
            except Exception as e:
                print(f"   ⚠️  Batch delete failed: {e}")

        print(f"\n✅ Deleted {deleted_count} files from index")
        return deleted_count

    except Exception as e:
        print(f"❌ Delete failed: {e}")
        return 0

def discover_files_in_folder(token, folder_id, folder_path, file_paths):
    """Recursively discover all file paths in a folder"""
    headers = {"Authorization": f"Bearer {token}"}
    base_url = "https://graph.microsoft.com/v1.0/me/drive"

    url = f"{base_url}/items/{folder_id}/children"
    response = requests.get(url, headers=headers)

    if response.status_code != 200:
        return

    items = response.json().get("value", [])

    # Add all files
    for item in items:
        if "file" in item:
            file_name = item['name']
            file_path = f"{folder_path}/{file_name}"
            file_paths.append(file_path)

    # Recurse into subfolders
    for item in items:
        if "folder" in item:
            subfolder_name = item['name']
            subfolder_path = f"{folder_path}/{subfolder_name}"
            subfolder_id = item['id']
            discover_files_in_folder(token, subfolder_id, subfolder_path, file_paths)

def load_folder_skip_cache():
    """Load folder skip preferences from cache"""
    if FOLDER_SKIP_CACHE_FILE.exists():
        with open(FOLDER_SKIP_CACHE_FILE, 'r') as f:
            return json.load(f)
    return {}

def save_folder_skip_cache(cache):
    """Save folder skip preferences to cache"""
    with open(FOLDER_SKIP_CACHE_FILE, 'w') as f:
        json.dump(cache, f, indent=2)

def should_process_folder(folder_path, folder_name, skip_cache, interactive=True):
    """Ask user if they want to process this folder
    Returns: 'process', 'skip', or 'list-only'
    """
    cache_key = folder_path

    # Check cache first
    if cache_key in skip_cache:
        decision = skip_cache[cache_key]
        if decision == "skip":
            print(f"📁 {folder_name}/ [SKIPPED - cached choice]")
            return "skip"
        elif decision == "process":
            print(f"📁 {folder_name}/ [PROCESSING - cached choice]")
            return "process"
        elif decision == "list-only":
            print(f"📁 {folder_name}/ [LIST ONLY - cached choice]")
            return "list-only"

    # Ask user if interactive
    if interactive:
        print(f"\n📁 Found folder: {folder_name}/")
        print(f"   Path: {folder_path}")
        while True:
            choice = input("   [y]es / [n]o / [l]ist-only / [a]lways yes / [s]kip always / [o]nly list always: ").lower().strip()
            if choice in ['y', 'yes', '']:
                return "process"
            elif choice in ['n', 'no']:
                return "skip"
            elif choice in ['l', 'list', 'list-only']:
                return "list-only"
            elif choice in ['a', 'always']:
                skip_cache[cache_key] = "process"
                save_folder_skip_cache(skip_cache)
                print(f"   ✅ Will always process this folder")
                return "process"
            elif choice in ['s', 'skip']:
                skip_cache[cache_key] = "skip"
                save_folder_skip_cache(skip_cache)
                print(f"   ⏭️ Will always skip this folder")
                return "skip"
            elif choice in ['o', 'only', 'only-list']:
                skip_cache[cache_key] = "list-only"
                save_folder_skip_cache(skip_cache)
                print(f"   📋 Will always list-only this folder")
                return "list-only"
            else:
                print("   Invalid choice, please try again.")

    # Default: process if not interactive
    return "process"

def list_folder_contents_only(token, folder_id, folder_path, extracted_files, processed_count):
    """Recursively list all files in a folder without extracting content (like ZIP list mode)"""
    headers = {"Authorization": f"Bearer {token}"}
    base_url = "https://graph.microsoft.com/v1.0/me/drive"

    # List items in current folder
    url = f"{base_url}/items/{folder_id}/children"
    response = requests.get(url, headers=headers)

    if response.status_code != 200:
        print(f"   ⚠️ Failed to access folder: {response.text}")
        return

    items = response.json().get("value", [])
    folders = [item for item in items if "folder" in item]
    files = [item for item in items if "file" in item]

    # Collect file listing
    file_list = []
    for item in files:
        file_name = item['name']
        file_size = item.get("size", 0)
        file_modified = item.get("lastModifiedDateTime", "")
        file_list.append(f"  • {file_name} ({file_size} bytes)")
        processed_count[0] += 1

    # Note: Subfolders are handled by crawl_folder_recursive(), not here
    # This function only lists files in the CURRENT folder

    # Create a single entry with the folder listing
    if file_list:
        folder_listing_text = f"=== Folder: {folder_path} ===\n"
        folder_listing_text += f"Contains {len(file_list)} files:\n"
        folder_listing_text += "\n".join(file_list)

        extracted_files.append({
            "name": f"{folder_path.split('/')[-1]} (folder listing)",
            "path": folder_path,
            "text": folder_listing_text,
            "size": 0,
            "modified": ""
        })
        print(f"   ✅ Listed {len(file_list)} files")

def is_token_expired(response):
    """Check if response indicates token expiration"""
    if response.status_code == 401:
        return True
    if response.status_code == 400:
        try:
            error_data = response.json()
            error_msg = str(error_data.get("error", {})).lower()
            return "invalidauthenticationtoken" in error_msg or "jwt" in error_msg or "token" in error_msg
        except:
            pass
    return False

def crawl_folder_recursive(token_ref, folder_id, folder_path, max_files, skip_cache, extracted_files, failed_files, skipped_files, processed_count, interactive=True):
    """Recursively crawl folders and extract files

    Args:
        token_ref: List containing the current access token [token]. Using a list allows updating the token reference.
    """
    headers = {"Authorization": f"Bearer {token_ref[0]}"}
    base_url = "https://graph.microsoft.com/v1.0/me/drive"

    # List items in current folder
    url = f"{base_url}/items/{folder_id}/children"
    response = requests.get(url, headers=headers)

    # If token expired, try to refresh it
    if is_token_expired(response):
        print(f"   🔄 Token expired, refreshing...")
        new_token = get_access_token(silent_only=True)
        if new_token:
            token_ref[0] = new_token  # Update token reference
            headers = {"Authorization": f"Bearer {token_ref[0]}"}
            response = requests.get(url, headers=headers)  # Retry request
            print(f"   ✅ Token refreshed successfully")
        else:
            print(f"   ❌ Token refresh failed")
            return processed_count[0]

    if response.status_code != 200:
        print(f"   ⚠️ Failed to access folder: {response.text}")
        return processed_count[0]

    items = response.json().get("value", [])

    # Process folders first (for interactive selection)
    folders = [item for item in items if "folder" in item]
    files = [item for item in items if "file" in item]

    # Check if CURRENT folder is list-only (skip file extraction if so)
    current_folder_mode = skip_cache.get(folder_path, "process")
    if current_folder_mode == "list-only":
        # List files only, don't extract content
        print(f"   📋 Listing files in {folder_path} (list-only mode)")
        list_folder_contents_only(
            token_ref[0], folder_id, folder_path, extracted_files, processed_count
        )
        return processed_count[0]

    # Process subfolders
    for folder_item in folders:
        if processed_count[0] >= max_files:
            break

        folder_name = folder_item['name']
        subfolder_path = f"{folder_path}/{folder_name}"
        subfolder_id = folder_item['id']

        # Ask if we should process this folder
        folder_mode = should_process_folder(subfolder_path, folder_name, skip_cache, interactive)

        if folder_mode == "process":
            # Full processing - recurse into subfolder
            crawl_folder_recursive(
                token_ref, subfolder_id, subfolder_path, max_files, skip_cache,
                extracted_files, failed_files, skipped_files, processed_count, interactive
            )
        elif folder_mode == "list-only":
            # List-only mode - just index filenames and paths (like ZIP list mode)
            print(f"   📋 Listing files only (not extracting content)")
            list_folder_contents_only(
                token_ref[0], subfolder_id, subfolder_path, extracted_files, processed_count
            )
        # else: skip - do nothing

    # Process files in current folder
    for item in files:
        if processed_count[0] >= max_files:
            print(f"\n⚠️ Reached {max_files} file limit")
            break

        file_name = item['name']
        file_ext = file_name.lower().split('.')[-1] if '.' in file_name else 'no extension'
        file_path_full = f"{folder_path}/{file_name}"
        print(f"📄 Processing: {file_path_full}")

        # Check if file already indexed in Pinecone (incremental sync - skip extraction!)
        file_modified = item.get("lastModifiedDateTime", "")
        file_size = item.get("size", 0)
        vector_id = generate_vector_id(file_path_full)

        try:
            result = index.fetch(ids=[vector_id], namespace="smartdrive")
            if vector_id in result.get('vectors', {}):
                # File exists in Pinecone - check if it's changed
                existing_vector = result['vectors'][vector_id]
                existing_modified = existing_vector.get('metadata', {}).get('modified', '')
                existing_size = existing_vector.get('metadata', {}).get('size', 0)

                if existing_modified == file_modified and existing_size == file_size:
                    # File unchanged - skip extraction entirely!
                    print(f"   ⏭️  Already indexed (unchanged) - skipping extraction")
                    processed_count[0] += 1
                    continue
                else:
                    # File changed - will re-extract
                    print(f"   🔄 File modified - re-indexing")
        except Exception as e:
            # File not in Pinecone or error checking - will extract
            pass

        text = extract_text_from_file(token_ref[0], item)

        if text:
            print(f"   ✅ Extracted {len(text)} characters")
            extracted_files.append({
                "name": file_name,
                "path": f"{folder_path}/{file_name}",
                "text": text,
                "size": item.get("size", 0),
                "modified": item.get("lastModifiedDateTime", "")
            })
        elif text is None:
            supported_extensions = ['pdf', 'docx', 'doc', 'pptx', 'xlsx', 'xlsm', 'xltx', 'xltm',
                                   'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif', 'txt', 'csv', 'json', 'md', 'markdown', 'zip']
            if file_ext in supported_extensions:
                failed_files.append((file_name, file_ext))
                print(f"   ❌ Failed to extract (see error above)")
            else:
                # Unsupported file type - index with metadata only
                skipped_files.append((file_name, file_ext))
                print(f"   📋 Unsupported type (.{file_ext}) - indexing metadata only")

                # Create metadata entry
                file_size_bytes = item.get("size", 0)
                file_size_mb = file_size_bytes / (1024 * 1024) if file_size_bytes > 0 else 0
                metadata_text = f"File: {file_name}\n"
                metadata_text += f"Type: .{file_ext} (unsupported for text extraction)\n"
                metadata_text += f"Size: {file_size_mb:.2f} MB ({file_size_bytes:,} bytes)\n"
                metadata_text += f"Location: {folder_path}/{file_name}"

                extracted_files.append({
                    "name": file_name,
                    "path": f"{folder_path}/{file_name}",
                    "text": metadata_text,
                    "size": item.get("size", 0),
                    "modified": item.get("lastModifiedDateTime", "")
                })

        processed_count[0] += 1

    return processed_count[0]

def discover_all_folders(token, folder_id, folder_path="/Documents", folders_list=None, failed_folders=None, depth=0):
    """Recursively discover all folders in the tree without processing files

    Returns: Tuple of (folders_list, failed_folders)
        folders_list: List of tuples [(folder_path, folder_name, folder_id), ...]
        failed_folders: List of tuples [(folder_path, folder_id, error_msg), ...]
    """
    if folders_list is None:
        folders_list = []
    if failed_folders is None:
        failed_folders = []

    headers = {"Authorization": f"Bearer {token}"}
    base_url = "https://graph.microsoft.com/v1.0/me/drive"

    # List items in current folder
    url = f"{base_url}/items/{folder_id}/children"

    try:
        response = requests.get(url, headers=headers, timeout=10)  # 10 second timeout
    except requests.exceptions.Timeout:
        print(f"   ⏱️ Timeout on {folder_path} (>10s), retrying with 30s timeout...")
        try:
            response = requests.get(url, headers=headers, timeout=30)  # Retry with longer timeout
        except Exception as retry_error:
            error_msg = f"Timeout after 30s"
            print(f"   ❌ Still timeout. Marking for later retry.")
            failed_folders.append((folder_path, folder_id, error_msg))
            return folders_list, failed_folders
    except requests.exceptions.RequestException as e:
        error_msg = f"Network error: {str(e)[:60]}"
        print(f"   ⚠️ {error_msg}")
        print(f"   ⚠️ Marking for later retry.")
        failed_folders.append((folder_path, folder_id, error_msg))
        return folders_list, failed_folders

    if response.status_code != 200:
        error_msg = f"API error {response.status_code}"
        print(f"   ⚠️ {error_msg} on {folder_path}. Marking for retry.")
        failed_folders.append((folder_path, folder_id, error_msg))
        return folders_list, failed_folders

    items = response.json().get("value", [])

    # Find all subfolders
    subfolder_items = [item for item in items if "folder" in item]

    for item in subfolder_items:
        folder_name = item['name']
        subfolder_path = f"{folder_path}/{folder_name}"
        subfolder_id = item['id']

        # Add this folder to the list
        folders_list.append((subfolder_path, folder_name, subfolder_id))

        # Show progress more frequently (every 5 folders)
        if len(folders_list) % 5 == 0:
            print(f"   ... {len(folders_list)} folders discovered")

        # Recursively discover subfolders
        discover_all_folders(token, subfolder_id, subfolder_path, folders_list, failed_folders, depth + 1)

    return folders_list, failed_folders

def interactive_folder_selection(folders_list, skip_cache):
    """Present all folders to user and let them choose how to handle each one

    Returns: Updated skip_cache with user preferences
    """
    print("\n" + "=" * 60)
    print("📁 FOLDER SELECTION - Choose how to handle each folder")
    print("=" * 60)
    print(f"Found {len(folders_list)} folders in your Documents directory.")
    print("\nFor each folder, choose:")
    print("  [y] = Process (extract all file contents)")
    print("  [l] = List-only (index filenames without extracting)")
    print("  [n] = Skip (ignore this folder)")
    print("  [x] = eXclude (skip this folder + ALL subfolders)")
    print("  [q] = Quick mode - assume 'yes' for all remaining folders")
    print("=" * 60)

    quick_mode = False
    excluded_prefixes = []  # Track folders that should auto-skip their subfolders
    included_prefixes = []  # Track folders that should auto-process their subfolders
    list_only_prefixes = []  # Track folders that should auto-list-only their subfolders

    for idx, (folder_path, folder_name, folder_id) in enumerate(folders_list, 1):
        # Check if this folder is under an excluded parent
        is_subfolder_of_excluded = any(folder_path.startswith(prefix + "/") for prefix in excluded_prefixes)
        if is_subfolder_of_excluded:
            skip_cache[folder_path] = "skip"
            print(f"\n[{idx}/{len(folders_list)}] {folder_name}/ → ⏭️  SKIP (parent excluded)")
            continue

        # Check if this folder is under a list-only parent
        is_subfolder_of_list_only = any(folder_path.startswith(prefix + "/") for prefix in list_only_prefixes)
        if is_subfolder_of_list_only:
            skip_cache[folder_path] = "list-only"
            print(f"\n[{idx}/{len(folders_list)}] {folder_name}/ → 📋 LIST-ONLY (parent list-only)")
            continue

        # Check if this folder is under an included parent
        is_subfolder_of_included = any(folder_path.startswith(prefix + "/") for prefix in included_prefixes)
        if is_subfolder_of_included:
            skip_cache[folder_path] = "process"
            print(f"\n[{idx}/{len(folders_list)}] {folder_name}/ → ✅ PROCESS (parent included)")
            continue

        # Check if we already have a cached decision
        if folder_path in skip_cache:
            decision = skip_cache[folder_path]
            status = {"process": "✅ PROCESS", "list-only": "📋 LIST-ONLY", "skip": "⏭️  SKIP"}[decision]
            print(f"\n[{idx}/{len(folders_list)}] {folder_name}/")
            print(f"   Path: {folder_path}")
            print(f"   {status} (cached)")
            continue

        # In quick mode, default to process
        if quick_mode:
            skip_cache[folder_path] = "process"
            print(f"\n[{idx}/{len(folders_list)}] {folder_name}/ → ✅ PROCESS (quick mode)")
            continue

        # Ask user
        print(f"\n[{idx}/{len(folders_list)}] {folder_name}/")
        print(f"   Path: {folder_path}")

        while True:
            choice = input("   Choice [y/l/n/x/q]: ").lower().strip()

            if choice in ['y', 'yes', '']:
                # Count how many subfolders will be auto-processed
                subfolder_count = sum(1 for p, _, _ in folders_list[idx:] if p.startswith(folder_path + "/"))
                skip_cache[folder_path] = "process"
                included_prefixes.append(folder_path)  # Mark for auto-processing subfolders
                if subfolder_count > 0:
                    print(f"   → ✅ Will PROCESS this folder + {subfolder_count} subfolders")
                else:
                    print("   → ✅ Will PROCESS this folder")
                break
            elif choice in ['l', 'list', 'list-only']:
                # Count how many subfolders will be list-only
                subfolder_count = sum(1 for p, _, _ in folders_list[idx:] if p.startswith(folder_path + "/"))
                skip_cache[folder_path] = "list-only"
                list_only_prefixes.append(folder_path)  # Mark for auto-list-only subfolders
                if subfolder_count > 0:
                    print(f"   → 📋 Will LIST-ONLY this folder + {subfolder_count} subfolders")
                else:
                    print("   → 📋 Will LIST-ONLY this folder")
                break
            elif choice in ['n', 'no', 'skip']:
                skip_cache[folder_path] = "skip"
                print("   → ⏭️  Will SKIP this folder")
                break
            elif choice in ['x', 'exclude']:
                # Count how many subfolders will be skipped
                subfolder_count = sum(1 for p, _, _ in folders_list[idx:] if p.startswith(folder_path + "/"))
                skip_cache[folder_path] = "skip"
                excluded_prefixes.append(folder_path)  # Mark for auto-skipping subfolders
                print(f"   → ⏭️  Will SKIP this folder + {subfolder_count} subfolders")
                break
            elif choice in ['q', 'quick']:
                skip_cache[folder_path] = "process"
                quick_mode = True
                print("   → ✅ Will PROCESS this folder")
                print("   🚀 Quick mode enabled - all remaining folders will be processed")
                break
            else:
                print("   Invalid choice, please enter y/l/n/x/q")

    # Save updated cache
    save_folder_skip_cache(skip_cache)

    print("\n" + "=" * 60)
    print("✅ Folder selection complete!")

    # Show summary
    process_count = sum(1 for v in skip_cache.values() if v == "process")
    list_count = sum(1 for v in skip_cache.values() if v == "list-only")
    skip_count = sum(1 for v in skip_cache.values() if v == "skip")

    print(f"   ✅ Process: {process_count} folders")
    print(f"   📋 List-only: {list_count} folders")
    print(f"   ⏭️  Skip: {skip_count} folders")
    print("=" * 60)

    return skip_cache

def list_documents_folder(token, max_files=None, interactive=True, preflight=True):
    """Recursively crawl Documents folder with interactive folder selection

    Args:
        token: OneDrive access token
        max_files: Maximum number of files to process (None for unlimited)
        interactive: If True, asks user about each folder (used when no cache)
        preflight: If True, discover all folders first and let user choose before crawling

    New Folder Detection:
        When using cached folder choices (interactive=False, preflight=False),
        automatically detects new folders not in cache and asks user about them
        before proceeding with the crawl. This gives you control over new folders
        while keeping speed benefits of cached choices for existing folders.
    """
    headers = {"Authorization": f"Bearer {token}"}
    base_url = "https://graph.microsoft.com/v1.0/me/drive"

    # Get Documents folder
    url = f"{base_url}/root:/Documents"
    response = requests.get(url, headers=headers)

    if response.status_code != 200:
        raise Exception(f"Failed to access Documents folder: {response.text}")

    folder_id = response.json()["id"]

    # Load folder skip cache
    skip_cache = load_folder_skip_cache()

    # NEW FOLDER DETECTION: If we have cached choices but NOT in full preflight mode,
    # check for new folders not in cache and ask about them
    if not interactive and not preflight and len(skip_cache) > 0:
        print(f"🔍 Checking for new folders not in cache...")
        folders_list, failed_folders = discover_all_folders(token, folder_id, "/Documents")

        # Find folders not in cache
        new_folders = [(path, name, fid) for path, name, fid in folders_list if path not in skip_cache]

        if new_folders:
            print(f"✨ Found {len(new_folders)} new folder(s) not in cache!\n")
            print("=" * 60)
            print("🆕 NEW FOLDERS DETECTED")
            print("=" * 60)
            print("How do you want to handle these new folders?")
            print("  [y] = Process (extract all file contents)")
            print("  [l] = List-only (index filenames without extracting)")
            print("  [n] = Skip (ignore this folder)")
            print("=" * 60 + "\n")

            # Ask about each new folder
            for folder_path, folder_name, folder_id_item in new_folders:
                print(f"📁 New folder: {folder_name}/")
                print(f"   Path: {folder_path}")

                while True:
                    choice = input("   Choice [y/l/n]: ").lower().strip()

                    if choice in ['y', 'yes', '']:
                        skip_cache[folder_path] = "process"
                        print(f"   → ✅ Will PROCESS this folder\n")
                        break
                    elif choice in ['l', 'list', 'list-only']:
                        skip_cache[folder_path] = "list-only"
                        print(f"   → 📋 Will LIST-ONLY this folder\n")
                        break
                    elif choice in ['n', 'no', 'skip']:
                        skip_cache[folder_path] = "skip"
                        print(f"   → ⏭️  Will SKIP this folder\n")
                        break
                    else:
                        print("   Invalid choice, please enter y/l/n")

            # Save updated cache
            save_folder_skip_cache(skip_cache)
            print("=" * 60)
            print("✅ New folder choices saved!")
            print("=" * 60 + "\n")

    # PREFLIGHT MODE: Discover all folders first, let user choose
    if preflight and interactive:
        print(f"🔍 Discovering all folders in Documents directory...")
        folders_list, failed_folders = discover_all_folders(token, folder_id, "/Documents")
        print(f"✅ Found {len(folders_list)} folders\n")

        # If some folders failed, offer retry
        retry_count = 0
        while failed_folders and retry_count < 3:
            print(f"\n⚠️  Warning: {len(failed_folders)} folder(s) failed during discovery")
            for folder_path, folder_id_failed, error_msg in failed_folders:
                print(f"   • {folder_path}: {error_msg}")

            retry_choice = input(f"\nRetry failed folders? [y/n]: ").lower().strip()
            if retry_choice not in ['y', 'yes']:
                print("⏭️  Skipping failed folders, continuing with discovered folders...")
                break

            # Retry failed folders
            print(f"\n🔄 Retrying {len(failed_folders)} failed folder(s)...")
            new_failed = []
            for folder_path, folder_id_failed, error_msg in failed_folders:
                print(f"   Retrying: {folder_path}...")
                partial_list, partial_failed = discover_all_folders(token, folder_id_failed, folder_path)
                # Add newly discovered folders to main list
                folders_list.extend(partial_list)
                # Track still-failed folders
                new_failed.extend(partial_failed)

            failed_folders = new_failed
            retry_count += 1

            if not failed_folders:
                print(f"✅ All folders recovered!\n")
            elif retry_count >= 3:
                print(f"⚠️  Still have {len(failed_folders)} failing folders after 3 retries")
                print(f"   Continuing with {len(folders_list)} successfully discovered folders\n")

        # Let user select how to handle each folder
        skip_cache = interactive_folder_selection(folders_list, skip_cache)

        print("\n🚀 Starting crawl with your folder preferences...\n")
        # Now interactive is False for the actual crawl (decisions already made)
        interactive = False

    # Initialize tracking
    extracted_files = []
    failed_files = []
    skipped_files = []
    processed_count = [0]  # Use list to allow modification in recursive calls

    # Set max_files to a very large number if None
    if max_files is None:
        max_files = 999999

    print(f"🔍 Crawling Documents folder recursively...")
    if interactive:
        print(f"💡 Tip: You'll be asked about each folder. Choose 'always' to remember your choice!\n")

    # Start recursive crawl (pass token as reference list for refresh capability)
    token_ref = [token]
    crawl_folder_recursive(
        token_ref, folder_id, "/Documents", max_files, skip_cache,
        extracted_files, failed_files, skipped_files, processed_count, interactive
    )

    # Print summary
    print(f"\n{'='*60}")
    print(f"📊 Processing Summary:")
    print(f"{'='*60}")
    print(f"✅ Successfully extracted: {len(extracted_files)} files")

    if failed_files:
        print(f"\n❌ Failed extractions ({len(failed_files)}):")
        for fname, ext in failed_files:
            print(f"   • {fname} (.{ext})")

    if skipped_files:
        print(f"\n⚠️ Unsupported file types ({len(skipped_files)}):")
        type_counts = {}
        for fname, ext in skipped_files:
            type_counts[ext] = type_counts.get(ext, 0) + 1
        for ext, count in type_counts.items():
            print(f"   • .{ext}: {count} file(s)")

    print(f"{'='*60}\n")

    return extracted_files

if __name__ == "__main__":
    import asyncio

    token = get_access_token()
    print("✅ Authentication successful!\n")

    # Check if folder cache exists
    skip_cache = load_folder_skip_cache()
    has_cache = len(skip_cache) > 0

    # Main menu
    while True:
        print("=" * 60)
        print("📋 SmartDrive Crawler - Main Menu")
        print("=" * 60)
        print("1. Run crawler (use cached folder choices)")
        print("2. Reset folder choices and start fresh")
        print("3. View/edit cached folder choices")
        print("4. Delete folder from index")
        print("5. Exit")
        print("=" * 60)

        if has_cache:
            print(f"ℹ️  You have {len(skip_cache)} cached folder choice(s)")
        else:
            print("ℹ️  No cached folder choices yet")

        choice = input("\nSelect option [1-5]: ").strip()

        if choice == "1":
            # Run crawler
            print("\n" + "=" * 60)

            # Ask about clearing the index
            print("🗑️  Index Management:")
            print("   - Press Enter to ADD to existing index (default)")
            print("   - Type 'clear' to CLEAR index before indexing (fresh start)")
            clear_choice = input("Index mode: ").strip().lower()

            if clear_choice == "clear":
                print("\n⚠️  WARNING: This will DELETE all existing vectors in the 'smartdrive' namespace!")
                confirm = input("Are you sure? Type 'yes' to confirm: ").strip().lower()
                if confirm == "yes":
                    print("🗑️  Clearing index...")
                    try:
                        index.delete(delete_all=True, namespace="smartdrive")
                        print("✅ Index cleared!\n")
                    except Exception as e:
                        # If namespace doesn't exist yet, that's fine
                        if "Namespace not found" in str(e) or "404" in str(e):
                            print("ℹ️  Namespace doesn't exist yet (this is normal for first run)\n")
                        else:
                            print(f"⚠️  Clear failed: {e}\n")
                else:
                    print("❌ Clear cancelled, will add to existing index\n")
            else:
                print("✅ Will add to existing index (default)\n")

            # Ask about ZIP file handling
            print("=" * 60)
            print("📦 ZIP File Handling:")
            print("   - Press Enter to LIST zip contents (default, faster)")
            print("   - Type 'extract' to EXTRACT and index files inside zips (slower)")
            zip_choice = input("ZIP handling: ").strip().lower()

            if zip_choice == "extract":
                EXTRACT_ZIP_CONTENTS = True
                print("✅ Will extract and index ZIP contents\n")
            else:
                EXTRACT_ZIP_CONTENTS = False
                print("✅ Will list ZIP contents only (default)\n")

            print("=" * 60)
            print("📋 File limit options:")
            print("   - Press Enter for NO LIMIT (process all files)")
            print("   - Or enter a number (e.g., 50 for testing)")
            limit_input = input("File limit: ").strip()

            if limit_input == "":
                max_files = None
                print("\n✅ No limit set - will process all files\n")
            else:
                try:
                    max_files = int(limit_input)
                    print(f"\n✅ Limit set to {max_files} files\n")
                except ValueError:
                    print("\n⚠️ Invalid input, defaulting to no limit\n")
                    max_files = None

            # If we have cached folder choices, skip the discovery phase
            if has_cache:
                files = list_documents_folder(token, max_files=max_files, interactive=False, preflight=False)
            else:
                files = list_documents_folder(token, max_files=max_files, interactive=True, preflight=True)
            break

        elif choice == "2":
            # Reset cache
            if has_cache:
                print(f"\n⚠️  This will delete {len(skip_cache)} cached folder choice(s).")
                confirm = input("Are you sure? [y/N]: ").strip().lower()
                if confirm in ['y', 'yes']:
                    if FOLDER_SKIP_CACHE_FILE.exists():
                        FOLDER_SKIP_CACHE_FILE.unlink()
                    print("✅ Folder choices reset!\n")
                    skip_cache = {}
                    has_cache = False
                else:
                    print("❌ Cancelled\n")
            else:
                print("\n⚠️  No cached folder choices to reset.\n")

        elif choice == "3":
            # View/edit cache
            if has_cache:
                print("\n" + "=" * 60)
                print("📂 Cached Folder Choices:")
                print("=" * 60)
                for idx, (folder_path, decision) in enumerate(skip_cache.items(), 1):
                    if decision == "process":
                        status = "✅ PROCESS"
                    elif decision == "skip":
                        status = "⏭️  SKIP"
                    else:
                        status = "📋 LIST-ONLY"
                    print(f"{idx}. {status} - {folder_path}")
                print("=" * 60)

                print("\nOptions:")
                print("  - Enter folder number to cycle through modes (process → list-only → skip → process)")
                print("  - Type 'delete #' to remove a cached choice (e.g., 'delete 3')")
                print("  - Press Enter to go back")

                edit_choice = input("\nYour choice: ").strip().lower()

                if edit_choice == "":
                    print()
                    continue
                elif edit_choice.startswith("delete "):
                    try:
                        num = int(edit_choice.split()[1])
                        folder_to_delete = list(skip_cache.keys())[num - 1]
                        del skip_cache[folder_to_delete]
                        save_folder_skip_cache(skip_cache)
                        print(f"✅ Removed cached choice for: {folder_to_delete}\n")
                        has_cache = len(skip_cache) > 0
                    except (IndexError, ValueError):
                        print("❌ Invalid number\n")
                else:
                    try:
                        num = int(edit_choice)
                        folder_to_toggle = list(skip_cache.keys())[num - 1]
                        current = skip_cache[folder_to_toggle]
                        # Cycle through: process -> list-only -> skip -> process
                        if current == "process":
                            new_value = "list-only"
                            status = "LIST-ONLY"
                        elif current == "list-only":
                            new_value = "skip"
                            status = "SKIP"
                        else:  # skip
                            new_value = "process"
                            status = "PROCESS"
                        skip_cache[folder_to_toggle] = new_value
                        save_folder_skip_cache(skip_cache)
                        print(f"✅ Changed to {status}: {folder_to_toggle}\n")
                    except (IndexError, ValueError):
                        print("❌ Invalid number\n")
            else:
                print("\n⚠️  No cached folder choices to view.\n")

        elif choice == "4":
            # Delete folder from index
            print("\n" + "=" * 60)
            print("🗑️  Delete Folder from Index")
            print("=" * 60)
            print("This will remove all files from a folder (and subfolders) from the index.")
            print("Example: /Documents/MyFolder")
            print()

            folder_path = input("Enter folder path to delete (or press Enter to cancel): ").strip()

            if folder_path == "":
                print("❌ Cancelled\n")
                continue

            # Confirm deletion
            print(f"\n⚠️  WARNING: This will DELETE all indexed files from:")
            print(f"   {folder_path}")
            print(f"   And ALL subfolders within it!")
            confirm = input("\nType 'delete' to confirm: ").strip().lower()

            if confirm == "delete":
                deleted_count = delete_folder_from_index(folder_path)

                # Also update folder cache to skip this folder in future
                if deleted_count > 0:
                    print(f"\n💡 Updating folder cache to skip this folder in future...")
                    skip_cache[folder_path] = "skip"
                    save_folder_skip_cache(skip_cache)
                    print(f"✅ Folder marked as 'skip' in cache\n")
            else:
                print("❌ Deletion cancelled\n")

        elif choice == "5":
            print("\n👋 Goodbye!\n")
            exit(0)

        else:
            print("\n❌ Invalid choice, please try again.\n")

    if files:
        upload_to_pinecone(files)

    # Clean up embedding provider sessions
    print("\n🧹 Cleaning up...")
    try:
        loop = asyncio.get_event_loop()
        loop.run_until_complete(embedding_provider.close())
    except Exception:
        pass
